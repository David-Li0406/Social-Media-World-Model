"""Generate a 4-slide summary deck of the new experiment results."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
GOOD = RGBColor(0x1E, 0x8B, 0x4E)
BAD = RGBColor(0xB0, 0x3A, 0x2E)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEC, 0xF2, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def setline(p, text, size=16, bold=False, color=None, align=None, bullet=False, space=4):
    p.text = ("• " + text) if bullet else text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Calibri"
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    p.space_after = Pt(space)


def header(s, title, sub=None):
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = box(s, 0.45, 0.12, 12.4, 1.0)
    setline(tf.paragraphs[0], title, size=27, bold=True, color=WHITE)
    if sub:
        setline(tf.add_paragraph(), sub, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))


def add_table(s, rows, l, t, w, h, col_w=None, header_fill=ACCENT, font=11,
              highlight_rows=None, highlight_color=LIGHT):
    highlight_rows = highlight_rows or {}
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_w:
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            c.margin_left = Pt(5); c.margin_right = Pt(5)
            p = c.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(font + (1 if i == 0 else 0))
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            if i == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = highlight_color if i in highlight_rows else WHITE
                if i in highlight_rows:
                    p.font.bold = True; p.font.color.rgb = NAVY
    return gt


# ---------------------------------------------------------------- Slide 1
s = slide()
bar = s.shapes.add_shape(1, 0, 0, SW, SH)
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
tf = box(s, 0.7, 1.5, 12, 2.2)
setline(tf.paragraphs[0], "Social Media World Model", size=40, bold=True, color=WHITE)
setline(tf.add_paragraph(), "Benchmark, Generalization & Reply-Summary Evaluation",
        size=22, color=RGBColor(0x9E, 0xC5, 0xE8))
tf = box(s, 0.7, 3.7, 12, 3)
setline(tf.paragraphs[0],
        "Task: predict a stimulus comment's engagement — score, controversiality, reply width, reply summary",
        size=16, color=WHITE, bullet=True)
setline(tf.add_paragraph(),
        "Headline metric: Spearman ρ (scale-invariant ranking) + pairwise accuracy; macro-F1 for controversiality",
        size=16, color=WHITE, bullet=True)
setline(tf.add_paragraph(),
        "What's new: 18 baselines (was 6) · cross-domain · temporal · data-scaling · LLM-as-judge on summaries",
        size=16, color=WHITE, bullet=True)
setline(tf.add_paragraph(),
        "Data: Reddit r/politics (7,596 train / 1,900 test) + 5 subreddits × 3 months for generalization",
        size=16, color=WHITE, bullet=True)
tf = box(s, 0.7, 6.7, 12, 0.5)
setline(tf.paragraphs[0], "Qwen3-4B LoRA fine-tuned on H20 GPUs via GitHub Actions runner",
        size=12, color=RGBColor(0x7E, 0xA5, 0xC8))

# ---------------------------------------------------------------- Slide 2
s = slide()
header(s, "1 · Main Benchmark (r/politics, 1,900 test)",
       "Cheap supervised models lead; conversation structure carries most of the signal")
rows = [
    ["Model", "score ρ", "width ρ", "contr F1", "type"],
    ["quantile_gbm (LightGBM)", "0.760", "0.687", "0.557", "tabular"],
    ["feature_gbdt", "0.759", "0.676", "0.424", "tabular"],
    ["structural_prior  (NO text)", "0.750", "0.670", "0.415", "graph-only"],
    ["gnn (reply-tree)", "0.745", "0.676", "0.494", "graph NN"],
    ["llm_reghead Qwen3-4B", "0.731", "0.663", "0.494", "LLM+heads"],
    ["llm_sft Qwen3-4B (JSON)", "0.724", "0.541", "0.545", "LLM SFT"],
    ["encoder (DistilBERT)", "0.720", "0.642", "0.494", "encoder"],
    ["Qwen3-32B zero-shot", "0.482", "0.477", "0.311", "LLM 0-shot"],
    ["Qwen3-4B zero-shot", "0.374", "0.508", "0.223", "LLM 0-shot"],
]
add_table(s, rows, 0.45, 1.4, 7.7, 4.9, col_w=[3.0, 1.2, 1.2, 1.2, 1.1],
          font=11.5, highlight_rows={1})
tf = box(s, 8.4, 1.45, 4.6, 5.2)
setline(tf.paragraphs[0], "Key findings", size=17, bold=True, color=NAVY, space=8)
pts = [
    ("Structure > text. ", "A text-free structural prior (ρ=0.750) and a reply-tree GNN (0.745) rival the best text model."),
    ("Cheap wins. ", "A sub-second LightGBM beats every LLM on ranking."),
    ("Reg-head > JSON. ", "llm_reghead lifts SFT 0.724→0.731 (score), 0.541→0.663 (width), 0% parse failures vs 81%."),
    ("Fine-tuning helps LLMs a lot. ", "Qwen3-4B 0.374→0.724; tuned 4B beats zero-shot 32B."),
]
for head, body in pts:
    p = tf.add_paragraph()
    p.text = "• " + head + body
    p.font.size = Pt(13); p.font.name = "Calibri"; p.space_after = Pt(9)
    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = ACCENT

# ---------------------------------------------------------------- Slide 3
s = slide()
header(s, "2 · Generalization: is it a world model?",
       "Trains on politics, transfers across domains & time; keeps improving with data")

# three figures (cross-domain / temporal / data-scaling) rendered by make_figures.py
fig_w = 4.28
for j, png in enumerate(["cross_domain.png", "temporal.png", "data_scaling.png"]):
    s.shapes.add_picture("slides/figures/" + png, Inches(0.18 + j * (fig_w + 0.05)),
                         Inches(1.35), width=Inches(fig_w))

tf = box(s, 0.45, 5.15, 12.6, 2.2)
setline(tf.paragraphs[0], "Takeaways", size=16, bold=True, color=NAVY, space=6)
for head, body in [
    ("Cross-domain: ", "graceful — only ≈0.02–0.09 ρ drop to 4 unseen subreddits; structural/GNN transfer best."),
    ("Temporal: ", "very stable — ≤0.03 ρ over 16 months back; engagement dynamics don't rot. Time-shift < domain-shift."),
    ("Scaling: ", "tree models monotonic & NOT saturated at 7,596 → more data should help; linear GLMs saturate early."),
]:
    p = tf.add_paragraph()
    p.text = "• " + head + body
    p.font.size = Pt(13.5); p.font.name = "Calibri"; p.space_after = Pt(7)
    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = GOOD

# ---------------------------------------------------------------- Slide 4
s = slide()
header(s, "3 · Reply-Summary Quality — LLM-as-Judge (Qwen3-32B)",
       "200 common records rated 1–5 vs the gold reference summary")
rows = [
    ["Model", "mean", "1", "2", "3", "4", "5"],
    ["Qwen3-4B zero-shot", "2.12", "11", "156", "32", "0", "1"],
    ["Qwen3-32B zero-shot", "2.07", "13", "161", "26", "0", "0"],
    ["llm_sft Qwen3-4B", "1.95", "33", "148", "16", "3", "0"],
]
add_table(s, rows, 0.45, 1.45, 6.6, 1.9,
          col_w=[2.7, 0.9, 0.6, 0.6, 0.6, 0.6, 0.6], font=12, highlight_rows={1})

tf = box(s, 0.45, 3.6, 6.7, 3.6)
setline(tf.paragraphs[0], "Findings", size=16, bold=True, color=NAVY, space=7)
for head, body, col in [
    ("Summary forecasting is unsolved. ", "All ~2/5 — predicting the content/sentiment of UNSEEN future replies is genuinely hard.", GREY),
    ("Fine-tuning HURT it. ", "SFT is worst (1.95) with the most 'unrelated' 1s — LoRA-SFT traded text quality for numeric/JSON.", BAD),
    ("Scale doesn't help. ", "32B ≈ 4B zero-shot.", GREY),
    ("Judge is reliable. ", "Discriminates 1/3/4–5 with specific reasons; 0% parse failures.", GREY),
]:
    p = tf.add_paragraph()
    p.text = "• " + head + body
    p.font.size = Pt(13); p.font.name = "Calibri"; p.space_after = Pt(8)
    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = col

# conclusions panel
panel = s.shapes.add_shape(1, Inches(7.45), Inches(1.45), Inches(5.45), Inches(5.6))
panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT; panel.line.color.rgb = ACCENT
tf = panel.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_top = Pt(12)
setline(tf.paragraphs[0], "Overall conclusions", size=17, bold=True, color=NAVY, space=10)
for t in [
    "Conversation STRUCTURE is the dominant, domain/time-transferable signal.",
    "A cheap LightGBM is the strongest, most scalable world model — beats fine-tuned & 32B LLMs on ranking.",
    "Numeric channel is strong (ρ≈0.76) and generalizes; TEXT (reply summary) channel is not yet usable.",
    "Numeric-vs-text trade-off: the SFT model best on numbers is the worst summarizer.",
    "Next: dedicated summary decoder; imbalance-aware controversiality; scale the dataset.",
]:
    p = tf.add_paragraph()
    p.text = "• " + t
    p.font.size = Pt(13); p.font.name = "Calibri"; p.space_after = Pt(9)
    p.font.color.rgb = NAVY

prs.save("slides/world_model_results.pptx")
print("saved slides/world_model_results.pptx")
